"""Build the CCA-F study kit as a PowerPoint deck (with embedded diagrams)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from content import (
    EXAM_FACTS, OFFICIAL_PREP_COURSES, CHAPTERS, DOMAINS,
    CHEAT_SHEET, QUESTIONS, OUT_OF_SCOPE, PRACTICAL_EXERCISES,
)

HERE = os.path.dirname(os.path.abspath(__file__))
DIAG = os.path.join(HERE, "diagrams")
OUT_PPTX = os.path.join(HERE, "CCA-Foundations-Study-Kit.pptx")

# Anthropic-ish palette
DARK = RGBColor(0x19, 0x19, 0x19)
ORANGE = RGBColor(0xD9, 0x77, 0x57)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
BLUE = RGBColor(0x3A, 0x6F, 0x8F)
GREEN = RGBColor(0x5A, 0x8A, 0x6B)
RED = RGBColor(0xB5, 0x48, 0x3D)
GREY = RGBColor(0x7A, 0x7A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_title_band(slide, prs, text, color=DARK, height=0.7):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  prs.slide_width, Inches(height))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()
    tf = band.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.1)
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CREAM


def add_textbox(slide, left, top, width, height, paragraphs,
               default_size=16, bullet=True, color=DARK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(paragraphs):
        if isinstance(item, tuple):
            text, size, bold, c = item
        else:
            text, size, bold, c = item, default_size, False, color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if bullet and not text.startswith("  "):
            p.text = "• " + text
        else:
            p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = c
        p.alignment = PP_ALIGN.LEFT
    return box


def add_image(slide, name, left, top, width):
    path = os.path.join(DIAG, name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                width=Inches(width))


def add_footer(slide, prs, text):
    box = slide.shapes.add_textbox(0, prs.slide_height - Inches(0.4),
                                  prs.slide_width, Inches(0.4))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = GREY


# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Cover
    s = add_blank_slide(prs)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                             prs.slide_width, prs.slide_height)
    band.fill.solid()
    band.fill.fore_color.rgb = DARK
    band.line.fill.background()
    band2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5),
                              prs.slide_width, Inches(2.0))
    band2.fill.solid()
    band2.fill.fore_color.rgb = ORANGE
    band2.line.fill.background()

    tb = s.shapes.add_textbox(0, Inches(2.6), prs.slide_width, Inches(2.0))
    tf = tb.text_frame
    tf.text = "CCA-F"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = CREAM
    p2 = tf.add_paragraph()
    p2.text = "Claude Certified Architect — Foundations"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(28)
    p2.font.color.rgb = CREAM

    tb2 = s.shapes.add_textbox(0, Inches(5.0), prs.slide_width, Inches(2.0))
    tf2 = tb2.text_frame
    tf2.text = "Exam-Ready Study Deck"
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = CREAM
    for line in [
        f"{EXAM_FACTS['questions']} questions · {EXAM_FACTS['duration_min']} minutes · pass = 720/1000",
        "5 weighted domains · 4 of 8 scenarios sampled per attempt",
    ]:
        p = tf2.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.color.rgb = CREAM

    # ---- Agenda
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Agenda")
    add_textbox(s, 0.6, 1.0, 12, 5.5, [
        "Exam overview & domain weighting",
        "Theory: 13 chapters across the Claude API, Agent SDK, MCP, Claude Code",
        "Domain-by-domain checklist (5 domains)",
        "The 8 production scenarios",
        "Cheat sheet — patterns the exam keeps testing",
        "Practice questions with answers and explanations",
        "Practical labs — what to actually build",
        "Recommended resources & 4-week roadmap",
        "Out-of-scope topics — don't waste study time here",
    ])
    add_footer(s, prs, "CCA-F Study Kit · Anthropic Skilljar prep + community guides")

    # ---- Exam overview
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Exam at a glance")
    add_textbox(s, 0.6, 0.9, 6, 5.5, [
        (f"Duration: {EXAM_FACTS['duration_min']} minutes", 18, True, BLUE),
        (f"Questions: {EXAM_FACTS['questions']}", 18, True, BLUE),
        (f"Format: {EXAM_FACTS['format']}", 16, False, DARK),
        (f"Scoring: {EXAM_FACTS['scaled_score']}", 16, False, DARK),
        (f"Scenarios sampled: {EXAM_FACTS['scenarios']}", 16, False, DARK),
        (f"Delivery: {EXAM_FACTS['delivery']}", 16, False, DARK),
        ("No guessing penalty — answer EVERY question", 16, True, RED),
    ], bullet=False)
    add_image(s, "01_domain_weights.png", left=6.8, top=0.9, width=6.2)

    # ---- Domain weighting (just diagram, big)
    s = add_blank_slide(prs)
    add_title_band(s, prs, "5 Domains, weighted")
    add_image(s, "01_domain_weights.png", left=2.5, top=1.0, width=8.5)

    # ---- Scenarios
    s = add_blank_slide(prs)
    add_title_band(s, prs, "The 8 production scenarios")
    add_textbox(s, 0.6, 1.0, 12, 5.5,
                EXAM_FACTS["scenarios_list"])

    # ---- Roadmap
    s = add_blank_slide(prs)
    add_title_band(s, prs, "4-week study roadmap")
    add_image(s, "14_roadmap.png", left=1.5, top=1.0, width=10.5)

    # ---- Theory chapters (one slide each)
    chapter_diagrams = {
        "1. Claude API — Fundamentals": None,
        "2. Tools and tool_use": "11_tool_choice.png",
        "3. Claude Agent SDK — Building Agentic Systems": "02_agent_loop.png",
        "4. Model Context Protocol (MCP)": "05_mcp_architecture.png",
        "5. Claude Code — Configuration & Workflows": "04_claude_md_hierarchy.png",
        "6. Prompt Engineering — Advanced Techniques": "12_few_shot.png",
        "7. Message Batches API": "07_sync_vs_batch.png",
        "8. Task Decomposition Strategies": None,
        "9. Escalation & Human-in-the-Loop": "13_escalation_triggers.png",
        "10. Error Handling in Multi-agent Systems": "08_error_categories.png",
        "11. Context Management in Production": "09_context_strategies.png",
        "12. Provenance Preservation in Multi-source Synthesis": None,
        "13. Claude Code Built-in Tools": None,
    }
    for title, points in CHAPTERS:
        s = add_blank_slide(prs)
        add_title_band(s, prs, title)
        diag = chapter_diagrams.get(title)
        if diag:
            add_textbox(s, 0.5, 1.0, 6.8, 5.8, points, default_size=14)
            add_image(s, diag, left=7.5, top=1.0, width=5.5)
        else:
            add_textbox(s, 0.5, 1.0, 12.3, 5.8, points, default_size=15)

    # ---- Hub-and-spoke deeper dive
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Multi-agent: Hub-and-Spoke")
    add_image(s, "03_hub_and_spoke.png", left=0.5, top=1.0, width=7.5)
    add_textbox(s, 8.2, 1.0, 4.8, 5.5, [
        ("Coordinator owns:", 16, True, ORANGE),
        ("  decomposition, delegation, aggregation, error handling", 14, False, DARK),
        ("Subagents:", 16, True, ORANGE),
        ("  isolated context — pass everything explicitly via Task prompt", 14, False, DARK),
        ("Parallelism:", 16, True, ORANGE),
        ("  multiple Task calls in one coordinator turn", 14, False, DARK),
        ("Watch for:", 16, True, RED),
        ("  overly narrow decomposition, generic error envelopes, silent suppression", 14, False, DARK),
    ], bullet=False)

    # ---- Hooks vs prompts
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Hooks vs Prompt Instructions")
    add_image(s, "06_hooks_vs_prompts.png", left=1.5, top=1.0, width=10.5)

    # ---- Planning vs direct
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Planning Mode vs Direct Execution")
    add_image(s, "10_planning_vs_direct.png", left=1.5, top=1.0, width=10.5)

    # ---- tool_choice modes
    s = add_blank_slide(prs)
    add_title_band(s, prs, "tool_choice — three modes")
    add_image(s, "11_tool_choice.png", left=1.5, top=1.0, width=10.5)

    # ---- Sync vs Batch
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Sync API vs Message Batches API")
    add_image(s, "07_sync_vs_batch.png", left=1.5, top=1.0, width=10.5)

    # ---- Error categories
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Error categories & responses")
    add_image(s, "08_error_categories.png", left=1.5, top=1.0, width=10.5)

    # ---- Context strategies
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Context-window management strategies")
    add_image(s, "09_context_strategies.png", left=1.5, top=1.0, width=10.5)

    # ---- Domain notes (one slide per domain)
    for d in DOMAINS:
        s = add_blank_slide(prs)
        add_title_band(s, prs, d["title"])
        flat = []
        for topic_title, points in d["topics"]:
            flat.append((topic_title, 16, True, BLUE))
            for p in points:
                flat.append(("  " + p, 13, False, DARK))
        add_textbox(s, 0.5, 0.95, 12.3, 6.3, flat, default_size=13, bullet=False)

    # ---- Cheat sheet (multi-slide)
    cs_chunks = [CHEAT_SHEET[i:i + 4] for i in range(0, len(CHEAT_SHEET), 4)]
    for chunk_index, chunk in enumerate(cs_chunks, start=1):
        s = add_blank_slide(prs)
        add_title_band(s, prs, f"Cheat sheet ({chunk_index}/{len(cs_chunks)})")
        flat = []
        for title, items in chunk:
            flat.append((title, 16, True, ORANGE))
            for it in items:
                flat.append(("  " + it, 12, False, DARK))
            flat.append(("", 12, False, DARK))
        add_textbox(s, 0.5, 0.95, 12.3, 6.3, flat, default_size=12, bullet=False)

    # ---- Practice questions (one slide per Q)
    for i, q in enumerate(QUESTIONS, start=1):
        s = add_blank_slide(prs)
        add_title_band(s, prs, f"Q{i} — {q['scenario']}")
        items = [(q["stem"], 14, True, DARK), ("", 12, False, DARK)]
        for c in q["choices"]:
            items.append((c, 13, False, DARK))
        items += [
            ("", 12, False, DARK),
            (f"Answer: {q['answer']}", 16, True, GREEN),
            ("Why: " + q["explain"], 12, False, DARK),
        ]
        add_textbox(s, 0.5, 0.95, 12.3, 6.3, items, default_size=12, bullet=False)

    # ---- Practical exercises
    for ex in PRACTICAL_EXERCISES:
        s = add_blank_slide(prs)
        add_title_band(s, prs, ex["title"])
        items = [(f"Domains: {ex['domains']}", 14, True, ORANGE), ("", 12, False, DARK)]
        for step in ex["steps"]:
            items.append((step, 14, False, DARK))
        add_textbox(s, 0.5, 0.95, 12.3, 6.3, items, default_size=14)

    # ---- Resources
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Skilljar prep tracks (free)")
    add_textbox(s, 0.5, 1.0, 12.3, 6.0, OFFICIAL_PREP_COURSES, default_size=15)

    s = add_blank_slide(prs)
    add_title_band(s, prs, "Documentation to study")
    docs = [
        "Claude API: Messages, Tool Use, Message Batches",
        "Agent SDK: Overview, Hooks, Subagents, Sessions",
        "Model Context Protocol: Tools, Resources, Servers",
        "Claude Code: Memory (CLAUDE.md), Skills, Hooks, Sub-agents, MCP, GitHub Actions, Headless",
        "Prompt Engineering Guide + Anthropic Cookbook",
    ]
    add_textbox(s, 0.5, 1.0, 12.3, 6.0, docs, default_size=16)

    s = add_blank_slide(prs)
    add_title_band(s, prs, "Community / video")
    add_textbox(s, 0.5, 1.0, 12.3, 6.0, [
        "YouTube playlist PLFz7SvAnfqLpjPCPJBkUy077BC5ecuE8c — exam walkthroughs and demos",
        "github.com/paullarionov/claude-certified-architect — full study guide & 76 practice Qs",
        "claudecertifications.com — free study guide + practice questions",
        "FlashGenius / Tutorialsdojo / DEV.to / Towards AI — domain primers and worked questions",
    ], default_size=15)

    # ---- Out of scope
    s = add_blank_slide(prs)
    add_title_band(s, prs, "Out of scope — don't waste time here", color=RED)
    add_textbox(s, 0.5, 1.0, 12.3, 6.0, OUT_OF_SCOPE, default_size=14)

    # ---- Closing
    s = add_blank_slide(prs)
    add_title_band(s, prs, "You are exam-ready when you can…")
    add_textbox(s, 0.5, 1.0, 12.3, 6.0, [
        ("…explain the agent loop end-to-end and the role of stop_reason", 16, True, BLUE),
        ("…choose between hook and prompt for any reliability requirement", 16, True, BLUE),
        ("…design 5 well-scoped MCP tools with structured error envelopes", 16, True, BLUE),
        ("…configure CLAUDE.md, .claude/rules, skills, .mcp.json correctly", 16, True, BLUE),
        ("…spot 'overly narrow decomposition' and 'attention dilution' anti-patterns", 16, True, BLUE),
        ("…know exactly when to choose Sync vs Batch", 16, True, BLUE),
        ("…enforce sequencing with preconditions and token-bound preview/execute pairs", 16, True, BLUE),
    ], default_size=16, bullet=False)

    prs.save(OUT_PPTX)
    print("PPTX built:", OUT_PPTX, os.path.getsize(OUT_PPTX), "bytes")


if __name__ == "__main__":
    build()
